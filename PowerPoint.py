from pptx import Presentation
from pptx.util import Inches
import os

# Define the path where the images are stored
image_path = "image_path"
# Define your variable names
variable_names = ["X0", "X1", "X2"]

# Create a presentation object
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]  # 0 is the index for a title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Your Presentation Title"  # Replace with your title
subtitle.text = "Subtitle or Additional Information"  # Replace with your subtitle

# Function to add a slide with three images and a title
def add_image_slide(prs, image_files, var_name):
    slide_layout = prs.slide_layouts[5]  # 5 is the index for a blank slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = var_name  # Set the title to the variable name
    
    # Calculate the total width of the slide and the width of each picture
    slide_width = prs.slide_width
    image_width = slide_width / 3
    left_margin = (slide_width - 3 * image_width) / 2  # Calculate the left margin to center the images

    top = Inches(2)
    for idx, image_file in enumerate(image_files):
        left_position = left_margin + idx * image_width  # Calculate the left position for each image
        slide.shapes.add_picture(image_file, left_position, top, width=image_width)

# Iterate over the variable names and add slides
for var in variable_names:
    image_files = [
        os.path.join(image_path, f"{var}.png"),
        os.path.join(image_path, f"{var}_1.png"),
        os.path.join(image_path, f"{var}_2.png")
    ]
    # Check if all image files exist before adding them to the slide
    if all(os.path.exists(image_file) for image_file in image_files):
        add_image_slide(prs, image_files, var)
    else:
        print(f"One or more images for {var} do not exist. Skipping this set.")

# Save the presentation
prs.save(r'path_save')


print("Presentation created successfully.")
